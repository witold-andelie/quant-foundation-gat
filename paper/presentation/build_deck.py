"""Build the professor-facing project presentation from main.pdf's content.

Generates three data-driven charts (matplotlib) and assembles a ~16-slide
PPTX matching the reference deck's Office theme. All numbers are taken from
paper/main.tex + paper/tables (authoritative), not from the earlier deck.
"""
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"
ASSETS.mkdir(exist_ok=True)

# ----------------------------------------------------------------------------
# Palette (Office theme, matched to the reference deck)
# ----------------------------------------------------------------------------
NAVY   = RGBColor(0x2A, 0x33, 0x4A)   # titles / dark text
BLUE   = RGBColor(0x44, 0x72, 0xC4)   # accent 1
LBLUE  = RGBColor(0x5B, 0x9B, 0xD5)   # accent 5
ORANGE = RGBColor(0xED, 0x7D, 0x31)   # warnings / energy
GREEN  = RGBColor(0x2E, 0x8B, 0x57)   # positive results
RED    = RGBColor(0xC0, 0x39, 0x2B)   # negative results
GREY   = RGBColor(0x7F, 0x7F, 0x7F)
LGREY  = RGBColor(0xBF, 0xBF, 0xBF)
DGREY  = RGBColor(0x59, 0x59, 0x59)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PANEL  = RGBColor(0xF2, 0xF4, 0xF8)   # light panel fill

HEXNAVY, HEXBLUE, HEXORANGE = "#2A334A", "#4472C4", "#ED7D31"
HEXGREEN, HEXRED, HEXGREY = "#2E8B57", "#C0392B", "#7F7F7F"

# ============================================================================
#  1. DATA-DRIVEN CHARTS
# ============================================================================
plt.rcParams.update({
    "font.family": "DejaVu Sans",
    "font.size": 15,
    "axes.edgecolor": "#B0B4BD",
    "axes.linewidth": 0.9,
    "figure.facecolor": "white",
    "axes.facecolor": "white",
})


def _clean(ax):
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)


def chart_equity_valueadd():
    """Attention value-add = GAT OOS Sharpe − uniform anchor's. Paper §5.7:
    positive in 30/30 seeded runs, spanning +0.69 to +3.03."""
    fig, ax = plt.subplots(figsize=(6.7, 4.0))
    lo, hi = 0.69, 3.03
    ax.barh([0], [hi - lo], left=lo, height=0.5, color=HEXBLUE, zorder=3)
    ax.axvline(0, color="#333333", linewidth=1.4, zorder=4)
    ax.set_xlim(-0.6, 3.7)
    ax.set_ylim(-0.9, 1.0)
    ax.set_yticks([])
    ax.set_xlabel("Attention value-add  (OOS Sharpe: GAT − uniform anchor)")
    ax.grid(axis="x", color="#E4E7EE", zorder=0)
    _clean(ax)
    ax.spines["left"].set_visible(False)
    ax.text((lo + hi) / 2, 0.0, "all 30 / 30 seeded runs", ha="center", va="center",
            color="white", fontweight="bold", fontsize=13.5)
    ax.text(lo, -0.45, "+0.69", ha="center", va="top", color=HEXBLUE, fontweight="bold", fontsize=13)
    ax.text(hi, -0.45, "+3.03", ha="center", va="top", color=HEXBLUE, fontweight="bold", fontsize=13)
    ax.text(0, 0.62, "0 = uniform anchor\n(no learned attention)", ha="center", va="bottom",
            color="#595959", fontsize=11)
    ax.set_title("Every seeded run lands above zero — learned attention always helps",
                 fontsize=13, color="#2A334A", pad=12)
    fig.tight_layout()
    fig.savefig(ASSETS / "eq_valueadd.png", dpi=200)
    plt.close(fig)


def chart_energy_clip():
    fig, ax = plt.subplots(figsize=(6.8, 4.0))
    groups = ["Trivial −price\npredictor", "GAT composite"]
    clipped = [6.28, 11.00]
    honest = [-1.52, -1.54]
    x = range(len(groups))
    w = 0.36
    b1 = ax.bar([i - w / 2 for i in x], clipped, w, label="As trained (±0.8 clip)",
                color=HEXORANGE, zorder=3)
    b2 = ax.bar([i + w / 2 for i in x], honest, w, label="Honest (plain return)",
                color=HEXNAVY, zorder=3)
    ax.axhline(0, color="#333333", linewidth=1.1, zorder=2)
    ax.set_xticks(list(x))
    ax.set_xticklabels(groups)
    ax.set_ylabel("OOS Sharpe ratio")
    ax.set_ylim(-3.2, 12.8)
    ax.grid(axis="y", color="#E4E7EE", zorder=0)
    _clean(ax)
    for bars in (b1, b2):
        for b in bars:
            v = b.get_height()
            off = 0.28 if v >= 0 else -0.28
            va = "bottom" if v >= 0 else "top"
            ax.text(b.get_x() + b.get_width() / 2, v + off, f"{v:+.2f}",
                    ha="center", va=va, fontsize=12.5, fontweight="bold",
                    color=(HEXORANGE if v >= 0 else HEXRED))
    ax.legend(frameon=False, fontsize=12, loc="upper left")
    ax.set_title("Remove the clip → the 'winner' loses money",
                 fontsize=14, color="#2A334A", pad=12)
    fig.tight_layout()
    fig.savefig(ASSETS / "en_clip.png", dpi=200)
    plt.close(fig)


def chart_energy_edge():
    fig, ax = plt.subplots(figsize=(6.8, 3.7))
    labels = ["edge persistence\n(spread today)", "edge ridge\n(both endpoints, no graph)", "edge GAT\n(message passing)"]
    vals = [0.000, 0.192, 0.248]
    colors = [HEXGREY, HEXGREY, HEXBLUE]
    bars = ax.barh(labels, vals, color=colors, height=0.6, zorder=3)
    ax.invert_yaxis()
    ax.set_xlabel("Forecast skill  (1 − MSE / MSE$_{persistence}$)")
    ax.set_xlim(0, 0.34)
    ax.grid(axis="x", color="#E4E7EE", zorder=0)
    _clean(ax)
    for b, v in zip(bars, vals):
        ax.text(v + 0.007, b.get_y() + b.get_height() / 2, f"{v:.3f}",
                va="center", ha="left", fontweight="bold", color="#2A334A", fontsize=13)
    ax.text(0.335, 0.45, "edge GAT beats the ridge\nby +0.056  (≈29%, 5/5 seeds)",
            fontsize=12, color=HEXGREEN, fontweight="bold", ha="right", va="center")
    fig.tight_layout()
    fig.savefig(ASSETS / "en_edge.png", dpi=200)
    plt.close(fig)


chart_equity_valueadd()
chart_energy_clip()
chart_energy_edge()
print("charts written")

# ============================================================================
#  2. PPTX ASSEMBLY
# ============================================================================
EMU_W, EMU_H = 12192000, 6858000
prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]

FOOTER = "Relational Alpha Factors with Graph Attention Networks  ·  Master's Capstone"


def _set_font(run, size, color, bold=False, italic=False, name="Calibri"):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = name


def box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def para(tf, runs, size=16, color=NAVY, bold=False, align=PP_ALIGN.LEFT,
         space_after=6, space_before=0, bullet=False, level=0, line=1.06, first=False):
    """runs: str or list of (text, color, bold[, size])."""
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    p.level = level
    try:
        p.line_spacing = line
    except Exception:
        pass
    if isinstance(runs, str):
        runs = [(runs, color, bold)]
    for item in runs:
        txt, col, bd = item[0], item[1], item[2]
        sz = item[3] if len(item) > 3 else size
        r = p.add_run()
        r.text = txt
        _set_font(r, sz, col, bd)
    if bullet:
        _add_bullet(p)
    return p


def _add_bullet(p, char="–", color=BLUE):
    pPr = p._pPr if p._pPr is not None else p.get_or_add_pPr()
    buClr = pPr.makeelement(qn("a:buClr"), {})
    srgb = buClr.makeelement(qn("a:srgbClr"), {"val": "%02X%02X%02X" % (color[0], color[1], color[2])})
    buClr.append(srgb)
    buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial"})
    buChar = pPr.makeelement(qn("a:buChar"), {"char": char})
    for el in (buClr, buFont, buChar):
        pPr.append(el)
    pPr.set("indent", str(-Inches(0.24).emu.__int__() if hasattr(Inches(0.24), "emu") else -219456))
    pPr.set("marL", "219456")


def rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w or 1)
    return sp


def header(slide, section, title, accent=BLUE):
    """Standard slide header: section label + title + accent rule."""
    _, tf = box(slide, 0.62, 0.34, 12.1, 0.5)
    para(tf, [(section.upper(), accent, True, 15)], first=True, space_after=0)
    _, tf2 = box(slide, 0.62, 0.72, 12.1, 0.9)
    para(tf2, [(title, NAVY, True, 27)], first=True, space_after=0, line=1.0)
    rect(slide, 0.64, 1.52, 1.5, 0.055, fill=accent)


def footer(slide, n):
    _, tf = box(slide, 0.62, 7.02, 9.5, 0.35)
    para(tf, [(FOOTER, LGREY, False, 10)], first=True, space_after=0)
    _, tf2 = box(slide, 11.6, 7.02, 1.1, 0.35)
    para(tf2, [(str(n), LGREY, False, 10)], first=True, space_after=0, align=PP_ALIGN.RIGHT)


def new(section=None, title=None, accent=BLUE, n=None):
    s = prs.slides.add_slide(BLANK)
    if section is not None:
        header(s, section, title, accent)
    footer(s, len(prs.slides._sldIdLst))  # auto page number = current slide count
    return s


def statcard(slide, l, t, w, h, big, label, accent=BLUE, big_size=32, sub=None):
    rect(slide, l, t, w, h, fill=PANEL)
    rect(slide, l, t, 0.09, h, fill=accent)
    _, tf = box(slide, l + 0.22, t + 0.14, w - 0.34, h - 0.24)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, [(big, accent, True, big_size)], first=True, space_after=2, line=1.0)
    para(tf, [(label, DGREY, False, 12.5)], space_after=0, line=1.02)
    if sub:
        para(tf, [(sub, GREY, False, 10.5)], space_after=0, line=1.0)


def pic_contain(slide, path, l, t, w, h, align_h="center", align_v="middle"):
    """Add image scaled to fit inside box (l,t,w,h) in inches, preserving aspect."""
    from PIL import Image
    iw, ih = Image.open(path).size
    ar = iw / ih
    bw, bh = w, h
    if bw / bh > ar:
        nh = bh
        nw = bh * ar
    else:
        nw = bw
        nh = bw / ar
    if align_h == "center":
        nl = l + (bw - nw) / 2
    elif align_h == "left":
        nl = l
    else:
        nl = l + (bw - nw)
    if align_v == "middle":
        nt = t + (bh - nh) / 2
    elif align_v == "top":
        nt = t
    else:
        nt = t + (bh - nh)
    slide.shapes.add_picture(str(path), Inches(nl), Inches(nt), Inches(nw), Inches(nh))


# ---------------------------------------------------------------------------
# SLIDE 1 — Title
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, fill=NAVY)
rect(s, 0, 0, 13.333, 0.28, fill=BLUE)
rect(s, 0, 5.02, 13.333, 0.055, fill=BLUE)
_, tf = box(s, 0.9, 0.95, 11.5, 0.6)
para(tf, [("CAPSTONE PROJECT  ·  MASTER'S REPORT", LBLUE, True, 16)], first=True)
_, tf = box(s, 0.9, 1.7, 11.6, 2.6)
para(tf, [("Relational Alpha Factors with", WHITE, True, 40)], first=True, space_after=2, line=1.02)
para(tf, [("Graph Attention Networks", WHITE, True, 40)], space_after=10, line=1.02)
para(tf, [("A dual-track study on US equities and European power markets", LBLUE, False, 21)], space_after=0, line=1.05)
_, tf = box(s, 0.9, 5.35, 11.5, 1.5)
para(tf, [("Wentao Ma", WHITE, True, 20)], first=True, space_after=4)
para(tf, [("One shared GAT kernel over two structurally different graphs  —  ",
           RGBColor(0xC7, 0xD0, 0xE0), False, 14),
          ("what does learned relational structure add?", LBLUE, False, 14)], space_after=4)
para(tf, [("github.com/witold-andelie/quant-foundation-gat   ·   July 2026",
           RGBColor(0x9A, 0xA6, 0xBD), False, 13)], space_after=0)

# ---------------------------------------------------------------------------
# SLIDE 2 — The big idea (motivation)
# ---------------------------------------------------------------------------
s = new("Motivation", "The core premise: assets are not islands", n=2)
_, tf = box(s, 0.62, 1.85, 6.35, 5.0)
para(tf, [("Classical quantitative signals score each asset ", NAVY, False, 17),
          ("in isolation", ORANGE, True, 17),
          (" — read one stock's own prices and volumes, emit one number.",
           NAVY, False, 17)], first=True, space_after=12)
para(tf, [("But the assets are coupled:", NAVY, True, 17)], space_after=8)
para(tf, [("Equities co-move within sectors and supply chains "
           "(statistical coupling).", DGREY, False, 15.5)], bullet=True, space_after=7, level=0)
para(tf, [("European power prices are literally wired together by "
           "cross-border transmission lines (physical coupling).", DGREY, False, 15.5)],
     bullet=True, space_after=14)
para(tf, [("The premise: ", NAVY, False, 17),
          ("this coupling structure is information", BLUE, True, 17),
          (".", NAVY, False, 17)], space_after=8)
para(tf, [("The project measures ", NAVY, False, 17),
          ("how much", NAVY, True, 17),
          (" — as a graph, where nodes are assets and edges are the "
           "couplings a model may exchange information along.", NAVY, False, 17)],
     space_after=0)
# right panel: island vs relational
rect(s, 7.35, 1.9, 5.35, 4.7, fill=PANEL)
_, tf = box(s, 7.65, 2.15, 4.8, 4.3)
para(tf, [("TWO FACTOR FAMILIES", GREY, True, 13)], first=True, space_after=12)
para(tf, [("Island factor", ORANGE, True, 18)], space_after=3)
para(tf, [("scores a node from its ", DGREY, False, 14),
          ("own history alone", DGREY, True, 14),
          (" (e.g. 5-day price reversal, volume shocks).", DGREY, False, 14)], space_after=16)
para(tf, [("Relational factor", BLUE, True, 18)], space_after=3)
para(tf, [("may also read the node's ", DGREY, False, 14),
          ("neighbours in a graph", DGREY, True, 14),
          (".", DGREY, False, 14)], space_after=18)
para(tf, [("The relational model's only privilege is ", NAVY, False, 14.5),
          ("topological", NAVY, True, 14.5),
          (":", NAVY, False, 14.5)], space_after=3, align=PP_ALIGN.CENTER)
para(tf, [("permission to exchange information along edges, with every other "
           "input held fixed. Each headline number is that margin.",
           GREY, False, 13)], align=PP_ALIGN.CENTER, space_after=0)

# ---------------------------------------------------------------------------
# SLIDE 3 — Primer: what is an alpha factor?
# ---------------------------------------------------------------------------
s = new("Primer  ·  for a non-finance audience", "What is an “alpha factor”?", n=3)
_, tf = box(s, 0.62, 1.8, 12.1, 1.4)
para(tf, [("An ", NAVY, False, 17),
          ("alpha factor", BLUE, True, 17),
          (" is a rule that assigns every asset a score at every point in time, "
           "built so that ", NAVY, False, 17),
          ("higher-scored assets should subsequently outperform lower-scored ones.",
           NAVY, True, 17)], first=True, space_after=8)
para(tf, [("The prediction is about ", NAVY, False, 17),
          ("relative", NAVY, True, 17),
          (" performance, not price levels. The object of interest is the ",
           NAVY, False, 17),
          ("cross-section", BLUE, True, 17),
          (" — all assets at one timestamp — and whether today's "
           "ranking of scores anticipates tomorrow's ranking of returns.",
           NAVY, False, 17)], space_after=0)
# three step panels
steps = [
    ("1", "Score the cross-section", "At each date, every asset gets a number from the factor rule.", BLUE),
    ("2", "Rank the scores", "Convert to a within-day ranking. Absolute magnitude is discarded.", LBLUE),
    ("3", "Check the ordering", "Did high-ranked assets out-return low-ranked ones over the next k days?", GREEN),
]
x = 0.62
for num, ttl, body, col in steps:
    w = 3.83
    rect(s, x, 3.5, w, 2.55, fill=PANEL)
    rect(s, x, 3.5, w, 0.62, fill=col)
    _, tf = box(s, x + 0.25, 3.58, w - 0.45, 0.5)
    para(tf, [(num + "   " + ttl, WHITE, True, 16)], first=True, space_after=0)
    _, tf = box(s, x + 0.25, 4.35, w - 0.5, 1.6)
    para(tf, [(body, DGREY, False, 14.5)], first=True, space_after=0, line=1.1)
    x += w + 0.2
_, tf = box(s, 0.62, 6.25, 12.1, 0.6)
para(tf, [("Individual factors are weak and correlated by design — the field builds "
           "them in large numbers and studies how to ", GREY, False, 14),
          ("combine", GREY, True, 14),
          (" them. That combination step is where this project's graph lives.",
           GREY, False, 14)], first=True, space_after=0)

# ---------------------------------------------------------------------------
# SLIDE 4 — Primer: scoring a factor
# ---------------------------------------------------------------------------
s = new("Primer  ·  for a non-finance audience", "How we grade a factor: two numbers", n=4)
# rank-IC card
rect(s, 0.62, 1.85, 5.9, 4.85, fill=PANEL)
rect(s, 0.62, 1.85, 5.9, 0.7, fill=BLUE)
_, tf = box(s, 0.9, 1.96, 5.4, 0.5)
para(tf, [("Rank-IC  (information coefficient)", WHITE, True, 18)], first=True, space_after=0)
_, tf = box(s, 0.95, 2.8, 5.25, 3.8)
para(tf, [("The (Spearman) correlation between the scores issued today and the "
           "returns realised afterwards, averaged over all dates.", DGREY, False, 15)],
     first=True, space_after=10)
para(tf, [("1.0", BLUE, True, 15), (" = perfect ordering", DGREY, False, 15)], bullet=True, space_after=4)
para(tf, [("0.0", BLUE, True, 15), (" = no predictive skill", DGREY, False, 15)], bullet=True, space_after=10)
para(tf, [("Reality check: ", NAVY, True, 15),
          ("real factors are weak.", DGREY, False, 15)], space_after=4)
para(tf, [("This project's best equity composite reaches rank-IC ≈ ",
           DGREY, False, 15),
          ("0.015", BLUE, True, 15),
          (" out-of-sample — small, but real and usable in combination.",
           DGREY, False, 15)], space_after=0)
# Sharpe card
rect(s, 6.8, 1.85, 5.9, 4.85, fill=PANEL)
rect(s, 6.8, 1.85, 5.9, 0.7, fill=NAVY)
_, tf = box(s, 7.08, 1.96, 5.4, 0.5)
para(tf, [("Sharpe ratio  (of a long–short book)", WHITE, True, 18)], first=True, space_after=0)
_, tf = box(s, 7.13, 2.8, 5.25, 3.8)
para(tf, [("Grade the factor as a strategy: ", DGREY, False, 15),
          ("buy", GREEN, True, 15), (" the top-ranked assets, ", DGREY, False, 15),
          ("short", RED, True, 15), (" the bottom-ranked ones, then divide mean "
           "return by its volatility (annualised).", DGREY, False, 15)], first=True, space_after=10)
para(tf, [("Long–short largely cancels overall market direction — what "
           "remains is the factor's own edge.", DGREY, False, 15)], bullet=True, space_after=10)
para(tf, [("Key subtlety: ", NAVY, True, 15),
          ("rank-IC and Sharpe can disagree.", DGREY, False, 15)], space_after=4)
para(tf, [("Ranking well says nothing about the ", DGREY, False, 15),
          ("magnitude", DGREY, True, 15),
          (" of wins and losses — a disagreement the energy track turns "
           "into a main finding.", DGREY, False, 15)], space_after=0)

# ---------------------------------------------------------------------------
# SLIDE 5 — The aim (answers the professor's recurring question directly)
# ---------------------------------------------------------------------------
s = new("The aim", "Why a graph model? What we are actually after", accent=BLUE)
# left panel — what it is NOT
lx, lw = 0.62, 5.98
rect(s, lx, 1.8, lw, 3.55, fill=PANEL)
rect(s, lx, 1.8, lw, 0.66, fill=GREY)
_, tf = box(s, lx + 0.28, 1.9, lw - 0.5, 0.5)
para(tf, [("What the GNN is NOT", WHITE, True, 17)], first=True, space_after=0)
_, tf = box(s, lx + 0.28, 2.64, lw - 0.55, 2.6)
para(tf, [("A bigger model bolted on to push prediction accuracy up.",
           DGREY, False, 14)], bullet=True, first=True, space_after=10)
para(tf, [("A way to feed ", DGREY, False, 14),
          ("extra data", DGREY, True, 14),
          (" into prediction — every exogenous driver (load, wind, solar, fuel) "
           "also goes to the graph-free baselines.", DGREY, False, 14)],
     bullet=True, space_after=10)
para(tf, [("A demo that “a graph beats a non-graph” — easy to show, hard to "
           "interpret.", DGREY, False, 14)], bullet=True, space_after=0)
# right panel — what it IS
rx, rw = 6.72, 5.98
rect(s, rx, 1.8, rw, 3.55, fill=PANEL)
rect(s, rx, 1.8, rw, 0.66, fill=BLUE)
_, tf = box(s, rx + 0.28, 1.9, rw - 0.5, 0.5)
para(tf, [("What the GNN IS", WHITE, True, 17)], first=True, space_after=0)
_, tf = box(s, rx + 0.28, 2.64, rw - 0.55, 2.6)
para(tf, [("A ", NAVY, False, 14),
          ("measurement instrument", BLUE, True, 14),
          (": it isolates the marginal value of relational structure — the "
           "coupling between assets — with the information set held fixed.",
           NAVY, False, 14)], bullet=True, first=True, space_after=10)
para(tf, [("The model's only privilege is ", NAVY, False, 14),
          ("topological", NAVY, True, 14),
          (": permission to exchange information along edges (correlations for "
           "equities, the physical grid for power).", NAVY, False, 14)],
     bullet=True, space_after=10)
para(tf, [("So every headline number is a ", NAVY, False, 14),
          ("margin", BLUE, True, 14),
          (" — the value of that permission, everything else equal.",
           NAVY, False, 14)], bullet=True, space_after=0)
# bottom band — the one-line aim
rect(s, 0.62, 5.62, 12.08, 1.02, fill=NAVY)
rect(s, 0.62, 5.62, 0.11, 1.02, fill=LBLUE)
_, tf = box(s, 1.0, 5.62, 11.5, 1.02)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
para(tf, [("The aim, in one line:   ", LBLUE, True, 15),
          ("measure what the coupling structure of markets is worth — and whether ",
           WHITE, False, 15),
          ("learning", WHITE, True, 15),
          (" to read it beats reading it ", WHITE, False, 15),
          ("naively", WHITE, True, 15),
          (".", WHITE, False, 15)], first=True, space_after=0, line=1.1)

# ---------------------------------------------------------------------------
# SLIDE 6 — Research questions (the aim, made concrete)
# ---------------------------------------------------------------------------
s = new("The question", "Three questions the aim breaks into", n=6)
_, tf = box(s, 0.62, 1.75, 12.1, 0.7)
para(tf, [("With the information set held fixed, the aim decomposes into three "
           "concrete, separately-measurable questions:", DGREY, False, 16.5)],
     first=True, space_after=0)
rqs = [
    ("RQ1", "Does learned attention add value over unlearned relational propagation, "
            "holding inputs and topology fixed?", BLUE),
    ("RQ2", "Does one GAT kernel transfer across structurally different graphs "
            "without manufacturing false positives?", LBLUE),
    ("RQ3", "On real power data, does the physical topology improve forecast skill, "
            "and where does its value concentrate?", ORANGE),
]
y = 2.9
for tag, body, col in rqs:
    rect(s, 0.62, y, 12.08, 0.95, fill=PANEL)
    rect(s, 0.62, y, 1.15, 0.95, fill=col)
    _, tf = box(s, 0.62, y, 1.15, 0.95)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, [(tag, WHITE, True, 20)], first=True, align=PP_ALIGN.CENTER, space_after=0)
    _, tf = box(s, 1.95, y, 10.55, 0.95)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, [(body, NAVY, False, 15.5)], first=True, space_after=0, line=1.05)
    y += 1.1

# ---------------------------------------------------------------------------
# SLIDE 6 — The attention A/B (central test)
# ---------------------------------------------------------------------------
s = new("Method  ·  the central test", "The attention A/B: two no-learning anchors", n=6)
pic_contain(s, ASSETS / "isolation_anchors.png", 0.55, 1.75, 7.55, 4.35, align_v="top")
_, tf = box(s, 8.25, 1.9, 4.5, 4.9)
para(tf, [("Every run carries ", NAVY, False, 15.5),
          ("three constructions", NAVY, True, 15.5),
          (" on one input matrix, scored by one harness:", NAVY, False, 15.5)],
     first=True, space_after=10)
para(tf, [("Island anchor", ORANGE, True, 14.5),
          (" — equal-weight mean, no graph.", DGREY, False, 14.5)], bullet=True, space_after=6)
para(tf, [("Uniform anchor", LBLUE, True, 14.5),
          (" — same mean, averaged over the graph, nothing learned.",
           DGREY, False, 14.5)], bullet=True, space_after=6)
para(tf, [("GAT", BLUE, True, 14.5),
          (" — learned attention on the same graph.", DGREY, False, 14.5)],
     bullet=True, space_after=12)
para(tf, [("island → uniform  =  ", DGREY, False, 14.5),
          ("what the graph adds", BLUE, True, 14.5)], space_after=4)
para(tf, [("uniform → GAT  =  ", DGREY, False, 14.5),
          ("what learned attention adds", BLUE, True, 14.5)], space_after=12)
para(tf, [("This isolates a ", NAVY, False, 14.5),
          ("mechanism", NAVY, True, 14.5),
          (", not a whole architecture — the price is smaller headline "
           "numbers.", NAVY, False, 14.5)], space_after=0)

# ---------------------------------------------------------------------------
# SLIDE 7 — Platform: how the GNN plugs in
# ---------------------------------------------------------------------------
s = new("Platform", "How the GNN plugs in — two seams, not a fork", n=7)
_, tf = box(s, 0.62, 1.75, 12.1, 0.8)
para(tf, [("The relational track joins an existing analytics platform "
           "(ingest → DuckDB warehouse → dbt marts → Streamlit) through ",
           DGREY, False, 15.5),
          ("two narrow interfaces", BLUE, True, 15.5),
          (", fixed in design records before implementation.", DGREY, False, 15.5)],
     first=True, space_after=0)
cards = [
    ("SEAM 1", "The propagator", BLUE,
     ["One snapshot's node features + topology in, one factor value per node out.",
      "Both the uniform baseline and the trained GAT are adapters on this one seam — so the A/B is a one-line swap, not two code paths.",
      "A characterisation test pins it: uniform propagation on a full graph reproduces the platform's existing cross-market mean exactly."]),
    ("SEAM 2", "The factor contract", LBLUE,
     ["One Factor contract over the (time, entity) panel absorbs both idioms.",
      "At the registry / apply / gate layer, island and relational factors are indistinguishable.",
      "That is what makes the four-gate comparison a fair one — every factor reaches evaluation through the same code."]),
]
x = 0.62
for tag, ttl, col, bullets in cards:
    w = 5.95
    rect(s, x, 2.75, w, 3.85, fill=PANEL)
    rect(s, x, 2.75, w, 0.72, fill=col)
    _, tf = box(s, x + 0.28, 2.84, w - 0.5, 0.55)
    para(tf, [(tag + "   ", WHITE, True, 13), (ttl, WHITE, True, 18)], first=True, space_after=0)
    _, tf = box(s, x + 0.28, 3.68, w - 0.55, 2.8)
    for i, b in enumerate(bullets):
        para(tf, [(b, DGREY, False, 13.5)], bullet=True, space_after=8, first=(i == 0), line=1.08)
    x += w + 0.2
_, tf = box(s, 0.62, 6.72, 12.1, 0.35)
para(tf, [("Research runs flow back into the warehouse as marts; 115 tests "
           "(55 in the capstone modules) pin the leakage-critical logic.",
           GREY, False, 12.5)], first=True, space_after=0)

# ---------------------------------------------------------------------------
# SLIDE 8 — Methodology: leakage-safe measurement
# ---------------------------------------------------------------------------
s = new("Method", "Leakage was treated as the #1 risk", n=8)
_, tf = box(s, 0.62, 1.75, 12.1, 0.55)
para(tf, [("Controls are ", DGREY, False, 15.5),
          ("structural", NAVY, True, 15.5),
          (" — asserted in code, not remembered by the experimenter. "
           "One split date is computed once and passed to every consumer.",
           DGREY, False, 15.5)], first=True, space_after=0)
# split protocol visual
segs = [("train", BLUE, 3.6), ("embargo", ORANGE, 0.9), ("valid", GREEN, 1.9),
        ("embargo", ORANGE, 0.9), ("out-of-sample", NAVY, 3.0)]
total = sum(w for _, _, w in segs)
x0, y0, W, H = 0.85, 2.75, 11.6, 0.85
x = x0
for name, col, w in segs:
    ww = W * w / total
    rect(s, x, y0, ww - 0.03, H, fill=col)
    _, tf = box(s, x, y0, ww - 0.03, H)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, [(name, WHITE, True, 13)], first=True, align=PP_ALIGN.CENTER, space_after=0)
    x += ww
_, tf = box(s, x0, y0 + H + 0.05, W, 0.4)
para(tf, [("Validation is embargoed on ", GREY, False, 12.5),
          ("both", GREY, True, 12.5),
          (" sides — labels reach k steps forward, so the trailing embargo "
           "(easy to forget) stops best-epoch selection leaking OOS returns.",
           GREY, False, 12.5)], first=True, space_after=0)
# three control cards
ctrls = [
    ("Negative control", "Shuffle each snapshot's labels → validation IC must collapse to zero. It does.", RED),
    ("Positive control", "Plant a recoverable signal → the same loop must recover it. It does (IC → 0.99).", GREEN),
    ("IC-loss objective", "Optimise −Pearson(pred, label) per snapshot; deploy the best-validation-epoch weights.", BLUE),
]
x = 0.62
for ttl, body, col in ctrls:
    w = 3.83
    rect(s, x, 4.35, w, 2.15, fill=PANEL)
    rect(s, x, 4.35, 0.09, 2.15, fill=col)
    _, tf = box(s, x + 0.25, 4.52, w - 0.45, 1.85)
    para(tf, [(ttl, col, True, 16)], first=True, space_after=6)
    para(tf, [(body, DGREY, False, 13.5)], space_after=0, line=1.1)
    x += w + 0.2

# ---------------------------------------------------------------------------
# SLIDE 9 — The four research gates
# ---------------------------------------------------------------------------
s = new("Method", "The four research gates", accent=BLUE, n=9)
_, tf = box(s, 0.62, 1.72, 12.1, 0.5)
para(tf, [("A relational composite is only interesting if it survives comparison "
           "against the island factors it was built from — all scored "
           "out-of-sample.", DGREY, False, 15.5)], first=True, space_after=0)
rows = [
    ("Value-added", "composite OOS Sharpe vs. best single alpha", "composite > max", "FAIL", RED),
    ("Consistency", "IS / OOS IC same sign, magnitude retained", "score ≥ 0.5  (0.63)", "PASS", GREEN),
    ("Uniqueness", "max |correlation| vs. each single alpha", "< 0.7  (0.26)", "PASS", GREEN),
    ("Robustness", "blend: OOS coverage, IC info-ratio, drawdown", "score ≥ 0.5  (0.68)", "PASS", GREEN),
]
y = 2.5
# header row
cols = [(0.62, 2.6), (3.22, 5.0), (8.22, 3.1), (11.32, 1.38)]
hd = ["GATE", "DEFINITION", "THRESHOLD  (equity)", "RESULT"]
rect(s, 0.62, y, 12.08, 0.55, fill=NAVY)
for (cx, cw), h in zip(cols, hd):
    _, tf = box(s, cx + 0.12, y, cw - 0.2, 0.55)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, [(h, WHITE, True, 13)], first=True, space_after=0,
         align=PP_ALIGN.CENTER if h == "RESULT" else PP_ALIGN.LEFT)
y += 0.55
for name, defn, thr, res, col in rows:
    rect(s, 0.62, y, 12.08, 0.82, fill=PANEL if res == "PASS" else RGBColor(0xFB, 0xEE, 0xEC))
    vals = [(name, NAVY, True, 15.5), (defn, DGREY, False, 14), (thr, DGREY, False, 14)]
    for (cx, cw), v in zip(cols[:3], vals):
        _, tf = box(s, cx + 0.12, y, cw - 0.2, 0.82)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        para(tf, [v], first=True, space_after=0, line=1.02)
    cx, cw = cols[3]
    _, tf = box(s, cx + 0.12, y, cw - 0.2, 0.82)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(tf, [(res, col, True, 15)], first=True, align=PP_ALIGN.CENTER, space_after=0)
    y += 0.82
_, tf = box(s, 0.62, y + 0.12, 12.1, 0.7)
para(tf, [("The composite is real, unique and consistent — it beats ",
           DGREY, False, 14.5),
          ("9 of its 10 inputs", GREEN, True, 14.5),
          (" — but never clears the strictest bar (beat the single best alpha). "
           "We report the open gate rather than soften it.", DGREY, False, 14.5)],
     first=True, space_after=0)

# ---------------------------------------------------------------------------
# SLIDE 10 — Equity headline result
# ---------------------------------------------------------------------------
s = new("Result  ·  equity track", "Learned attention reliably beats naive propagation", accent=BLUE, n=10)
statcard(s, 0.62, 1.8, 3.85, 1.55, "30 / 30", "seeded runs with positive attention value-add",
         accent=GREEN, big_size=34, sub="2 HP configs · 2 devices · +0.69 to +3.03 Sharpe")
statcard(s, 0.62, 3.5, 3.85, 1.55, "+0.0148", "best-config OOS IC  (± 0.0043)",
         accent=BLUE, big_size=32, sub="all 5 seeds positive · mean 3.4σ above zero")
statcard(s, 0.62, 5.2, 3.85, 1.5, "3 / 4", "research gates pass on real data",
         accent=BLUE, big_size=32, sub="Value-added is the one honest failure")
pic_contain(s, ASSETS / "eq_valueadd.png", 4.75, 2.0, 4.55, 4.6, align_v="top")
_, tf = box(s, 9.5, 1.95, 3.25, 4.9)
para(tf, [("What the number means", NAVY, True, 16)], first=True, space_after=8)
para(tf, [("The value-add is the GAT's OOS Sharpe ", DGREY, False, 13.5),
          ("minus the uniform anchor's", DGREY, True, 13.5),
          (" — same inputs, same topology, differing only in whether the edge "
           "weights are learned.", DGREY, False, 13.5)], space_after=8)
para(tf, [("Both no-learning anchors are ", DGREY, False, 13.5),
          ("negative", RED, True, 13.5),
          (" out-of-sample; the learned model is consistently positive.",
           DGREY, False, 13.5)], space_after=8)
para(tf, [("It is ", DGREY, False, 13.5),
          ("not learning to smooth", BLUE, True, 13.5),
          (": its scores are near-uncorrelated with the uniform anchor "
           "(Spearman −0.29..−0.14) — it learns ", DGREY, False, 13.5),
          ("where", DGREY, True, 13.5),
          (" the graph matters.", DGREY, False, 13.5)], space_after=8)
para(tf, [("Data: real yfinance, 49 names, 2021–2026, 10 island alphas, "
           "correlation-top-k + sector graph.", GREY, False, 11.5)], space_after=0)

# ---------------------------------------------------------------------------
# SLIDE 11 — What the attention learned
# ---------------------------------------------------------------------------
s = new("Result  ·  equity track", "What did the attention actually learn?", accent=BLUE, n=11)
pic_contain(s, ASSETS / "att_neighbour_weight.png", 0.55, 4.0, 8.0, 2.85, align_v="top")
_, tf = box(s, 0.62, 1.75, 7.9, 2.1)
para(tf, [("Attention analysis on 919,336 edge-weight rows over 1,364 dates:",
           NAVY, True, 15)], first=True, space_after=7)
para(tf, [("Genuinely relational", GREEN, True, 14),
          (" — self-loop weight only 0.084, so ", DGREY, False, 14),
          ("91.6% of attention mass sits on neighbours", DGREY, True, 14),
          (" (the graph is not decoration).", DGREY, False, 14)], bullet=True, space_after=5)
para(tf, [("But nearly uniform", ORANGE, True, 14),
          (" — across ~12.5 neighbours, entropy 0.956, sector-homophily lift "
           "just +0.019. A ", DGREY, False, 14),
          ("gentle re-weighting of mean pooling", DGREY, True, 14),
          (", not a sharp selector.", DGREY, False, 14)], bullet=True, space_after=0)
_, tf = box(s, 8.75, 1.9, 4.0, 4.9)
para(tf, [("The mechanism, in one line", NAVY, True, 15.5)], first=True, space_after=8)
para(tf, [("Small learned tilts on top of broad pooling — which is exactly "
           "why the GAT's edge is ", DGREY, False, 14),
          ("consistent yet modest", BLUE, True, 14),
          (".", DGREY, False, 14)], space_after=10)
para(tf, [("Temporally stationary", GREEN, True, 14),
          (": the neighbour share holds at 0.91–0.92 across 2021–2026, "
           "with no break at the IS/OOS line (chart).", DGREY, False, 14)],
     space_after=10)
para(tf, [("The design hypothesis — attention adapts to macro regimes — "
           "is ", DGREY, False, 14),
          ("not supported", RED, True, 14),
          (" at this level, and we say so.", DGREY, False, 14)], space_after=10)
para(tf, [("Validity check: the biggest attention sinks are mega-cap "
           "bellwethers (AMZN, NVDA, MSFT, GS…) — the names the universe "
           "co-moves with.", GREY, False, 12.5)], space_after=0)

# ---------------------------------------------------------------------------
# SLIDE 12 — Energy: the physical graph
# ---------------------------------------------------------------------------
s = new("Track 2  ·  energy", "The energy graph is physical, not estimated", accent=ORANGE, n=12)
pic_contain(s, ASSETS / "interconnector.png", 0.5, 1.75, 7.3, 5.1, align_v="top")
_, tf = box(s, 8.0, 2.0, 4.75, 4.8)
para(tf, [("The plan: repeat the equity experiment on a ", NAVY, False, 15.5),
          ("physically grounded", ORANGE, True, 15.5),
          (" graph.", NAVY, False, 15.5)], first=True, space_after=10)
para(tf, [("20 European bidding zones", NAVY, True, 14.5),
          (" (day-ahead spot prices from ENTSO-E).", DGREY, False, 14.5)],
     bullet=True, space_after=6)
para(tf, [("38 physical borders", NAVY, True, 14.5),
          (" — land interconnectors and HVDC cables.", DGREY, False, 14.5)],
     bullet=True, space_after=6)
para(tf, [("Nothing is estimated", BLUE, True, 14.5),
          (" — the topology is the grid itself.", DGREY, False, 14.5)],
     bullet=True, space_after=6)
para(tf, [("DE_LU is the natural hub", NAVY, True, 14.5),
          (" (11 interconnectors) — the physical layout, not a choice.",
           DGREY, False, 14.5)], bullet=True, space_after=12)
para(tf, [("Same GAT kernel, training loop, gates and A/B as the equity track "
           "— only the graph, label horizon (24h) and node set differ.",
           GREY, False, 13)], space_after=0)

# ---------------------------------------------------------------------------
# SLIDE 13 — Energy cautionary tale
# ---------------------------------------------------------------------------
s = new("Track 2  ·  energy", "A worked cautionary tale: three nested artifacts", accent=ORANGE, n=13)
pic_contain(s, ASSETS / "en_clip.png", 0.55, 1.8, 6.2, 4.6, align_v="top")
_, tf = box(s, 7.0, 1.8, 5.7, 5.0)
para(tf, [("The first real-data run looked spectacular: OOS IC 0.20, ",
           DGREY, False, 14.5),
          ("Sharpe 8.25, 4/4 gates", RED, True, 14.5),
          (". The same kernel earns ~1.4 on equities — ", DGREY, False, 14.5),
          ("a symptom, not a discovery.", NAVY, True, 14.5)], first=True, space_after=9)
para(tf, [("Overlapping labels", ORANGE, True, 14),
          (" — 24h returns sampled hourly share 23/24 hours (autocorr 0.91); "
           "Sharpe treats them as independent.", DGREY, False, 14)], bullet=True, space_after=5)
para(tf, [("Wrong annualisation", ORANGE, True, 14),
          (" — the equity constant (252) applied to hourly data (8,760).",
           DGREY, False, 14)], bullet=True, space_after=5)
para(tf, [("A clip on the evaluation return", ORANGE, True, 14),
          (" — a ±0.8 clip capped exactly the short-leg scarcity spikes "
           "where the money is lost.", DGREY, False, 14)], bullet=True, space_after=9)
para(tf, [("Under honest returns the strategy ", NAVY, False, 14.5),
          ("loses money", RED, True, 14.5),
          (" (Sharpe ≈ −1.5) — yet rank-IC stays ≈ 0.21.",
           DGREY, False, 14.5)], space_after=5)
para(tf, [("Lesson: never transform the evaluation return, and never read rank "
           "quality as profitability.", NAVY, True, 13.5)], space_after=0)

# ---------------------------------------------------------------------------
# SLIDE 14 — Energy reframe: forecasting
# ---------------------------------------------------------------------------
s = new("Track 2  ·  energy reframe", "Does the physical graph improve forecasts?", accent=ORANGE, n=14)
pic_contain(s, ASSETS / "energy_ladder.png", 0.55, 1.8, 5.2, 5.0, align_v="top")
_, tf = box(s, 6.0, 1.85, 6.7, 5.0)
para(tf, [("Reframed from “is there alpha?” to “does the interconnector "
           "graph improve forecast skill?” — a skill ladder where each rung "
           "adds ", DGREY, False, 15),
          ("exactly one ingredient", NAVY, True, 15),
          (".", DGREY, False, 15)], first=True, space_after=10)
statcard(s, 6.0, 3.35, 6.55, 1.35, "+0.131",
         "skill from the graph alone (uniform neighbour mean vs. no-graph ridge)",
         accent=GREEN, big_size=30,
         sub="pure numpy, no learned model — the hardest number to argue with; "
             "synthetic uncoupled zones show no such lift")
para(tf, [("", NAVY, False, 4)], space_after=0)
_, tf2 = box(s, 6.0, 4.95, 6.6, 1.9)
para(tf2, [("Attention vs. uniform", BLUE, True, 14),
           (" — wins on ", DGREY, False, 14),
           ("ranking", DGREY, True, 14),
           (" (5/5 seeds) but is a wash on MSE skill: uniform pooling is already "
            "near-optimal here.", DGREY, False, 14)], first=True, bullet=True, space_after=6)
para(tf2, [("Congestion as an edge feature", ORANGE, True, 14),
           (" — an honest ", DGREY, False, 14),
           ("null", RED, True, 14),
           (" under two operationalisations and two GAT backends.", DGREY, False, 14)],
     bullet=True, space_after=0)

# ---------------------------------------------------------------------------
# SLIDE 15 — Energy relational payoff (edge)
# ---------------------------------------------------------------------------
s = new("Track 2  ·  energy reframe", "The relational payoff: cross-border spreads", accent=ORANGE, n=15)
pic_contain(s, ASSETS / "en_edge.png", 0.5, 1.85, 6.6, 4.6, align_v="top")
_, tf = box(s, 7.3, 1.9, 5.4, 5.0)
para(tf, [("A cross-border spread (price$_a$ − price$_b$) is ",
           NAVY, False, 15.5),
          ("irreducibly relational", BLUE, True, 15.5),
          (" — it does not exist for a single node, and it is what "
           "transmission rights actually price.", NAVY, False, 15.5)],
     first=True, space_after=10)
para(tf, [("The GAT message-passes over the grid; a small head predicts each "
           "border's next-day spread.", DGREY, False, 14.5)], bullet=True, space_after=6)
para(tf, [("The baseline already sees ", DGREY, False, 14.5),
          ("both", DGREY, True, 14.5),
          (" endpoints' drivers + the current spread — so any gain is "
           "whole-network context, not missing data.", DGREY, False, 14.5)],
     bullet=True, space_after=10)
para(tf, [("+0.056 skill (≈29%), 5/5 seeds", GREEN, True, 15.5),
          (" — and the synthetic control ", DGREY, False, 14.5),
          ("inverts", NAVY, True, 14.5),
          (": on independent zones the GAT is worse than the ridge.",
           DGREY, False, 14.5)], space_after=8)
para(tf, [("The GNN's value concentrates precisely on the target that cannot be "
           "defined without the relation.", NAVY, True, 14.5)], space_after=0)

# ---------------------------------------------------------------------------
# SLIDE 16 — Conclusion
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, fill=NAVY)
rect(s, 0, 0, 13.333, 0.28, fill=BLUE)
_, tf = box(s, 0.7, 0.55, 12, 0.9)
para(tf, [("CONCLUSION", LBLUE, True, 15)], first=True, space_after=2)
para(tf, [("Three questions, answered directly", WHITE, True, 28)], space_after=0)
ans = [
    ("RQ1  Learned attention?", "Yes on equities — positive value-add in 30/30 runs against a "
     "matched anchor. The mechanism is a modest, stationary tilt on broad pooling; the composite "
     "beats 9 of 10 inputs but not the best single alpha.", GREEN),
    ("RQ2  One kernel, two graphs?", "Yes — the same code finds a real margin on equities and "
     "correctly finds nothing on synthetic random walks and uncoupled zones. Implausible magnitudes "
     "were treated as alarms, not results.", LBLUE),
    ("RQ3  Physical topology helps forecasts?", "Yes, where it should — +0.131 node skill, and "
     "+0.056 on cross-border spreads (5/5), while a synthetic control inverts. Congestion edge "
     "features: an honest null.", ORANGE),
]
y = 1.7
for ttl, body, col in ans:
    rect(s, 0.7, y, 11.95, 1.15, fill=RGBColor(0x36, 0x40, 0x59))
    rect(s, 0.7, y, 0.11, 1.15, fill=col)
    _, tf = box(s, 1.0, y + 0.13, 11.5, 0.95)
    para(tf, [(ttl + "   ", col, True, 15.5),
              (body, RGBColor(0xE6, 0xEA, 0xF2), False, 13.5)], first=True, space_after=0, line=1.06)
    y += 1.28
_, tf = box(s, 0.7, y + 0.05, 11.95, 1.2)
para(tf, [("The most transferable output is methodological: ",
           WHITE, True, 14.5),
          ("evaluate on untransformed returns; treat rank correlation and "
           "profitability as different claims; distrust magnitudes that outrun "
           "their plausibility; require seeds, fixed devices and "
           "cross-implementation replication; and keep negative controls running "
           "— they are what give positive results their meaning.",
           RGBColor(0xC7, 0xD0, 0xE0), False, 14.5)], first=True, space_after=0, line=1.12)

out = HERE / "GAT_Capstone_Presentation.pptx"
prs.save(out)
print("saved", out, "with", len(prs.slides._sldIdLst), "slides")
